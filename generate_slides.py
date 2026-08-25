"""
generate_slides.py
把 reports/ 底下任何一份報告 Markdown（月報或年報）自動轉成同檔名的 .pptx 簡報，
之後跑 export_static_data.py 時會自動一併複製、讓網站上的「下載簡報」按鈕生效。

用法：
    python generate_slides.py reports/2025-annual.md
    python generate_slides.py reports/2026-07.md

會在同一個資料夾產生對應的 .pptx（例如 reports/2025-annual.pptx）。

視覺設計跟網站一致（深色底、鏡頭藍／燈泡琥珀雙色系），每個項目符號會保留
Markdown 裡 **粗體** 標出的關鍵詞（用琥珀色加粗顯示），項目太多時自動分頁，
避免文字擠出版面。

【圖表語法】在任何一個 section 底下，可以用下面這種 fenced code block 標記
數據，腳本會自動畫成真正的 PowerPoint 長條圖（不是圖片，PowerPoint 裡可以
直接點兩下編輯數據）：

    ```chart
    title: 出貨量比較（萬台）
    series: 2024, 2025
    全球: 2017.9, 1920.6
    中國智能投影: 604.3, 520.3
    ```

- 第一行 title 是圖表標題
- 第二行 series 是每個類別要比較的欄位（例如年份），用逗號分隔
- 之後每一行是「類別名稱: 數值1, 數值2, ...」，數值順序要對應 series
- 一個 section 裡可以放多個 chart 區塊，每個都會各自產生一張圖表投影片

仍然是全自動產生、不需要另外裝 Node.js；如果要更講究、含時間軸的版本，
可以請 AI 助手用 pptxgenjs 客製化製作。

【版面自動判斷】如果一個段落底下的項目符號，大多數是「**關鍵詞**：說明」這種
格式、而且數量不太多（2-8 條），會自動改用卡片網格呈現（比較像簡報，不是條列）；
不符合這個格式、或項目太多太長的段落，維持原本的條列式呈現，並沿用動態分頁機制。
"""
import re
import sys
import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_LABEL_POSITION

# 跟網站一致的配色（見 style.css / tailwind 設計 token）
ROOM = RGBColor(0x14, 0x16, 0x1A)
CARD = RGBColor(0x1C, 0x1F, 0x26)
CARD2 = RGBColor(0x23, 0x26, 0x2F)
SCREEN = RGBColor(0xF5, 0xF3, 0xEC)
MIST = RGBColor(0x90, 0x96, 0xA1)
LENS = RGBColor(0x4C, 0x7E, 0xFF)
LAMP = RGBColor(0xFF, 0xB4, 0x54)
PURPLE = RGBColor(0x8B, 0x7F, 0xD6)
GREEN = RGBColor(0x6F, 0xCF, 0x97)
BORDER = RGBColor(0x2A, 0x2E, 0x38)

# 卡片邊框顏色輪替色盤（跟架構簡報那份的配色邏輯一致）
CARD_PALETTE = [LENS, LAMP, PURPLE, GREEN]

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

CHARS_PER_LINE = 46          # 保守估計（寧可低估，讓分頁提早，也不要讓文字溢出）
LINE_HEIGHT_IN = 0.26        # 每行估計高度（含行距）
BULLET_GAP_IN = 0.18         # 每條項目符號之間的間距
CONTENT_TOP_IN = 1.75
CONTENT_BOTTOM_IN = 6.85     # 底線留給頁碼/邊界

# 卡片網格：每頁最多放幾張卡片（太多會太擠，改分頁）
MAX_CARDS_PER_SLIDE = 6

# 段落標題關鍵字 → 圖示，抓不到對應關鍵字就不顯示圖示（不會影響版面）
SECTION_ICONS = [
    (("市場數據", "市場回顧", "銷量", "出貨"), "📊"),
    (("新品", "品牌", "發布"), "✨"),
    (("技術", "供應鏈"), "🔧"),
    (("關注", "展望", "趨勢"), "🔭"),
    (("摘要", "重點"), "📝"),
    (("時間軸", "Timeline", "轉折"), "🕐"),
]


def _section_icon(heading: str) -> str:
    for keywords, icon in SECTION_ICONS:
        if any(k in heading for k in keywords):
            return icon
    return ""


def _runs_char_len(runs):
    return sum(len(text) for text, _ in runs)


def _estimate_bullet_height(runs):
    char_len = max(1, _runs_char_len(runs))
    lines = max(1, -(-char_len // CHARS_PER_LINE))  # 無條件進位
    return lines * LINE_HEIGHT_IN + BULLET_GAP_IN


def paginate_bullets(bullets):
    """依動態估算高度把 bullets 切成好幾頁，每頁塞得下多少算多少，而不是固定條數。"""
    if not bullets:
        return [[]]

    budget = CONTENT_BOTTOM_IN - CONTENT_TOP_IN
    pages, current, used = [], [], 0.0

    for runs in bullets:
        h = _estimate_bullet_height(runs)
        if current and used + h > budget:
            pages.append(current)
            current, used = [], 0.0
        current.append(runs)
        used += h

    if current:
        pages.append(current)
    return pages


def _extract_card_term(runs, max_term_len=18):
    """
    判斷這條項目符號是不是「**關鍵詞**：說明」格式，是的話回傳 (term, description)，
    不是的話回傳 None。關鍵詞長度設上限，避免整句話剛好被寫成粗體時被誤判成卡片標題。
    """
    if not runs or not runs[0][1]:  # 第一段不是粗體，不符合格式
        return None
    term = runs[0][0].strip()
    term = re.sub(r"[：:]\s*$", "", term)  # 去掉可能包含在粗體裡的冒號
    if not term or len(term) > max_term_len:
        return None
    description = "".join(text for text, _ in runs[1:]).strip()
    description = re.sub(r"^[：:]\s*", "", description)  # 去掉開頭的冒號
    if not description:
        return None
    return term, description


def _is_card_friendly(bullets, min_ratio=0.7):
    """段落裡大多數項目都符合「關鍵詞：說明」格式、且數量不太多，才適合改用卡片網格。"""
    if not bullets or len(bullets) > 12:
        return False
    matched = sum(1 for b in bullets if _extract_card_term(b) is not None)
    return (matched / len(bullets)) >= min_ratio


def paginate_cards(bullets, per_page=MAX_CARDS_PER_SLIDE):
    """把符合卡片格式的項目，依 per_page 切成好幾頁。"""
    cards = [_extract_card_term(b) or (b[0][0][:18], "".join(t for t, _ in b[1:])) for b in bullets]
    return [cards[i:i + per_page] for i in range(0, len(cards), per_page)] or [[]]


# ---------------------------------------------------------------------------
# 解析報告 Markdown
# ---------------------------------------------------------------------------

def _split_bold_runs(text: str):
    """把 '**關鍵詞**：其餘文字' 拆成 [(text, is_bold), ...]，供簡報上分色顯示。"""
    parts = []
    for chunk in re.split(r"(\*\*.+?\*\*)", text):
        if not chunk:
            continue
        if chunk.startswith("**") and chunk.endswith("**"):
            parts.append((chunk[2:-2], True))
        else:
            parts.append((chunk, False))
    return parts


BULLET_PREFIX = re.compile(r"^[-*•]\s+")


def _parse_chart_block(lines: list[str]) -> dict | None:
    """把 ```chart ... ``` 區塊內的文字，解析成 {"title":..., "series":[...], "categories":[(label,[values]), ...]}。"""
    chart = {"title": "", "series": [], "categories": []}
    for line in lines:
        line = line.strip()
        if not line:
            continue
        if line.lower().startswith("title:"):
            chart["title"] = line.split(":", 1)[1].strip()
        elif line.lower().startswith("series:"):
            chart["series"] = [s.strip() for s in line.split(":", 1)[1].split(",") if s.strip()]
        elif ":" in line:
            label, raw_values = line.split(":", 1)
            try:
                values = [float(v.strip()) for v in raw_values.split(",") if v.strip()]
            except ValueError:
                continue  # 格式不對就跳過這一行，不讓整份報告解析失敗
            if values:
                chart["categories"].append((label.strip(), values))

    if not chart["categories"]:
        return None  # 沒有任何一列有效數據，這個圖表區塊當作沒寫
    if not chart["series"]:
        # 沒寫 series 的話，用「數值1」「數值2」...當預設欄位名稱
        n = len(chart["categories"][0][1])
        chart["series"] = [f"數值{i+1}" for i in range(n)]
    return chart


def parse_report(md_text: str):
    """
    把報告 Markdown 拆成 {title, sections: [{heading, bullets: [...], charts: [...]}]}。

    相容兩種常見寫法：
    - 條列式：「- 內容」或「*   內容」（dash 或 asterisk 開頭，空格數量不拘）
    - 整段文字：標題底下直接寫一段話，沒有條列符號（例如月報的「本月摘要」）
    兩種都會被當成一個個「項目」放進 bullets，簡報上統一用圓點呈現。

    另外支援 ```chart ... ``` fenced code block，解析成 charts（見檔案開頭的語法說明），
    畫成真正的 PowerPoint 圖表。
    """
    lines = md_text.splitlines()
    title = ""
    sections = []
    current = None
    in_chart = False
    chart_buffer: list[str] = []

    for raw_line in lines:
        line = raw_line.strip()

        if in_chart:
            if line == "```":
                in_chart = False
                if current is not None:
                    chart = _parse_chart_block(chart_buffer)
                    if chart:
                        current["charts"].append(chart)
                chart_buffer = []
            else:
                chart_buffer.append(raw_line)
            continue

        if not line:
            continue

        if line == "```chart":
            in_chart = True
            chart_buffer = []
            continue

        if line.startswith("# ") and not title:
            title = line[2:].strip()
            continue

        if line.startswith("## "):
            current = {"heading": line[3:].strip(), "bullets": [], "charts": []}
            sections.append(current)
            continue

        if line.startswith("---"):
            continue

        if current is None:
            continue  # 標題出現前的內容（理論上不會有）不處理

        text = BULLET_PREFIX.sub("", line)  # 有條列符號就去掉，沒有就原樣使用
        # 拿掉 [1][2] 這種引用標記，簡報上不需要，完整引用留在網頁版報告
        text = re.sub(r"\s*(?:\[\d+\])+", "", text)
        if text:
            current["bullets"].append(_split_bold_runs(text))

    return {"title": title or "報告", "sections": sections}


# ---------------------------------------------------------------------------
# 底層繪圖輔助
# ---------------------------------------------------------------------------

def _set_background(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_shape_alpha(shape, alpha_pct):
    """python-pptx 沒有直接支援填色透明度，用底層 XML 補上 <a:alpha>。"""
    sp_pr = shape.fill._xPr
    solid_fill = sp_pr.find(qn("a:solidFill"))
    if solid_fill is None:
        return
    srgb = solid_fill.find(qn("a:srgbClr"))
    if srgb is None:
        return
    alpha = srgb.makeelement(qn("a:alpha"), {"val": str(int(alpha_pct * 1000))})
    srgb.append(alpha)


def _add_glow_circle(slide, cx, cy, d, color, alpha_pct):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, d, d)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    _set_shape_alpha(shape, alpha_pct)
    return shape


def _add_textbox(slide, x, y, w, h, text, size, color, bold=False,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Arial"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return box


def _add_page_number(slide, n):
    _add_textbox(slide, Inches(12.5), Inches(7.05), Inches(0.6), Inches(0.3),
                 str(n).zfill(2), 10, MIST, align=PP_ALIGN.RIGHT)


def _add_card_rect(slide, x, y, w, h, line_color=BORDER, fill_color=CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = 0.06
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    return shape


def _add_eyebrow_heading(slide, section_no, heading):
    """段落編號 + 圖示 + 標題，內容頁跟卡片頁共用同一套視覺。"""
    icon = _section_icon(heading)
    eyebrow_text = f"{str(section_no).zfill(2)} · SECTION" + (f"　{icon}" if icon else "")
    _add_textbox(slide, Inches(0.6), Inches(0.42), Inches(8), Inches(0.35), eyebrow_text, 11, LAMP, bold=True)
    _add_textbox(slide, Inches(0.6), Inches(0.78), Inches(12), Inches(0.75), heading, 26, SCREEN, bold=True)


# ---------------------------------------------------------------------------
# 投影片版面
# ---------------------------------------------------------------------------

def build_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, ROOM)

    # 呼應網站首頁的光束/光暈意象
    _add_glow_circle(slide, Inches(9.3), Inches(-2.4), Inches(7), LENS, 14)
    _add_glow_circle(slide, Inches(-2.6), Inches(4.6), Inches(6), LAMP, 12)

    _add_textbox(slide, Inches(0.8), Inches(2.15), Inches(10), Inches(0.4),
                 "投影機情報站", 13, LAMP, bold=True)
    _add_textbox(slide, Inches(0.8), Inches(2.6), Inches(11.2), Inches(2.0),
                 title, 40, SCREEN, bold=True)
    _add_textbox(slide, Inches(0.8), Inches(4.55), Inches(10.5), Inches(0.5),
                 subtitle, 15, MIST)
    return slide


def build_content_slide(prs, section_no, heading, bullet_group, continued=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, ROOM)
    _add_eyebrow_heading(slide, section_no, heading)

    if not bullet_group:
        _add_textbox(slide, Inches(0.6), Inches(1.9), Inches(11), Inches(0.5), "（本段無內容）", 14, MIST)
        _add_page_number(slide, section_no)
        return slide

    y = CONTENT_TOP_IN
    for runs in bullet_group:
        h = _estimate_bullet_height(runs)
        y_in = Inches(y)

        # 圓點項目符號
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.65), y_in + Inches(0.14), Inches(0.11), Inches(0.11))
        dot.fill.solid()
        dot.fill.fore_color.rgb = LENS
        dot.line.fill.background()
        dot.shadow.inherit = False

        box = slide.shapes.add_textbox(Inches(0.95), y_in, Inches(11.7), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.line_spacing = 1.15
        for text, is_bold in runs:
            if not text:
                continue
            run = p.add_run()
            run.text = text
            run.font.size = Pt(14.5)
            run.font.name = "Arial"
            if is_bold:
                run.font.bold = True
                run.font.color.rgb = LAMP
            else:
                run.font.color.rgb = SCREEN
        y += h

    _add_page_number(slide, section_no)
    return slide


def build_card_grid_slide(prs, section_no, heading, cards):
    """
    卡片網格版面，取代單調的條列，用在「關鍵詞：說明」格式的內容上
    （例如品牌動態、新品盤點這類一條一個主題的段落）。
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, ROOM)
    _add_eyebrow_heading(slide, section_no, heading)

    n = len(cards)
    if n <= 2:
        cols = max(n, 1)
    elif n == 3:
        cols = 3
    elif n == 4:
        cols = 4
    else:
        cols = 3
    rows = -(-n // cols)  # 無條件進位

    margin_x = 0.6
    gap = 0.25
    top = 1.75
    bottom_margin = 0.55
    grid_w = 13.333 - margin_x * 2
    grid_h = 7.5 - top - bottom_margin

    card_w = (grid_w - gap * (cols - 1)) / cols
    card_h = (grid_h - gap * (rows - 1)) / rows

    for idx, (term, desc) in enumerate(cards):
        r, c = divmod(idx, cols)
        x = Inches(margin_x + c * (card_w + gap))
        y = Inches(top + r * (card_h + gap))
        w = Inches(card_w)
        h = Inches(card_h)

        accent = CARD_PALETTE[idx % len(CARD_PALETTE)]
        _add_card_rect(slide, x, y, w, h, line_color=accent, fill_color=CARD2)

        pad = Inches(0.22)
        _add_textbox(slide, x + pad, y + Inches(0.18), w - pad * 2, Inches(0.5),
                     term, 14.5, accent, bold=True)

        desc_box = slide.shapes.add_textbox(x + pad, y + Inches(0.72), w - pad * 2, h - Inches(0.9))
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.line_spacing = 1.2
        run = p.add_run()
        run.text = desc
        run.font.size = Pt(11.5)
        run.font.name = "Arial"
        run.font.color.rgb = SCREEN

    _add_page_number(slide, section_no)
    return slide


def build_chart_slide(prs, section_no, heading, chart):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, ROOM)
    _add_eyebrow_heading(slide, section_no, heading)

    chart_data = CategoryChartData()
    chart_data.categories = [label for label, _ in chart["categories"]]
    for i, series_name in enumerate(chart["series"]):
        chart_data.add_series(series_name, [values[i] for _, values in chart["categories"]])

    x, y, cx, cy = Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.2)
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    )
    gchart = graphic_frame.chart

    if chart.get("title"):
        gchart.has_title = True
        gchart.chart_title.text_frame.text = chart["title"]
        title_run = gchart.chart_title.text_frame.paragraphs[0].runs[0]
        title_run.font.size = Pt(15)
        title_run.font.bold = True
        title_run.font.color.rgb = SCREEN
        title_run.font.name = "Arial"
    else:
        gchart.has_title = False

    # 圖例
    gchart.has_legend = len(chart["series"]) > 1
    if gchart.has_legend:
        gchart.legend.position = XL_LEGEND_POSITION.BOTTOM
        gchart.legend.include_in_layout = False
        gchart.legend.font.size = Pt(11)
        gchart.legend.font.color.rgb = MIST
        gchart.legend.font.name = "Arial"

    # 數列顏色：跟網站配色一致（鏡頭藍／燈泡琥珀，超過兩個就用中性灰接續）
    palette = [LENS, LAMP, MIST]
    plot = gchart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(10)
    plot.data_labels.font.color.rgb = SCREEN
    plot.data_labels.font.name = "Arial"
    for i, series in enumerate(plot.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = palette[i % len(palette)]

    # 座標軸文字顏色調成配合深色底
    try:
        cat_axis = gchart.category_axis
        cat_axis.tick_labels.font.size = Pt(11)
        cat_axis.tick_labels.font.color.rgb = MIST
        cat_axis.tick_labels.font.name = "Arial"
        cat_axis.format.line.color.rgb = BORDER

        val_axis = gchart.value_axis
        val_axis.tick_labels.font.size = Pt(10)
        val_axis.tick_labels.font.color.rgb = MIST
        val_axis.format.line.color.rgb = BORDER
        val_axis.major_gridlines.format.line.color.rgb = BORDER
    except Exception:
        pass  # 座標軸樣式屬於錦上添花，失敗也不影響圖表本身資料正確性

    _add_page_number(slide, section_no)
    return slide


def build_closing_slide(prs, page_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, ROOM)
    _add_glow_circle(slide, Inches(9.8), Inches(3.5), Inches(6), LENS, 10)
    _add_textbox(slide, Inches(0.8), Inches(3.2), Inches(10), Inches(0.9),
                 "投影機情報站", 30, SCREEN, bold=True)
    _add_textbox(slide, Inches(0.8), Inches(3.95), Inches(10), Inches(0.5),
                 "由本站爬蟲每日追蹤產業動態，AI 自動摘要與分類", 14, MIST)
    _add_page_number(slide, page_no)
    return slide


# ---------------------------------------------------------------------------
# 主流程
# ---------------------------------------------------------------------------

def generate(md_path: str):
    with open(md_path, "r", encoding="utf-8") as f:
        report = parse_report(f.read())

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_title_slide(prs, report["title"], "投影機產業報告 · 自動產生簡報")

    page = 2
    for i, section in enumerate(report["sections"], start=1):
        has_chart = bool(section.get("charts"))
        for chart in section.get("charts", []):
            build_chart_slide(prs, page, section["heading"], chart)
            page += 1

        bullets = section["bullets"]
        if bullets and _is_card_friendly(bullets):
            for group in paginate_cards(bullets):
                build_card_grid_slide(prs, page, section["heading"], group)
                page += 1
        elif bullets or not has_chart:
            groups = paginate_bullets(bullets)
            for g_idx, group in enumerate(groups):
                build_content_slide(prs, page, section["heading"], group, continued=(g_idx > 0))
                page += 1

    build_closing_slide(prs, page)

    out_path = os.path.splitext(md_path)[0] + ".pptx"
    prs.save(out_path)
    print(f"簡報已產出：{out_path}（共 {page} 頁）")


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("用法：python generate_slides.py <報告的 .md 路徑>")
        sys.exit(1)
    generate(sys.argv[1])
